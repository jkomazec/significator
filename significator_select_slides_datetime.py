import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import os
import re
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from io import BytesIO
from pptx.enum.text import PP_ALIGN
from datetime import datetime

timestamp = datetime.now().strftime('%Y-%m-%d-%H-%M')

st.title('Dapresy export significator')

st.subheader('Import your pptx/ppt file')

#da mogu da odaberu font
font_list = ['Arial','Roboto','Calibri', 'Times New Roman', 'Consolas']
selected_font = st.selectbox('Select  font: ', font_list)
#da mogu da odaberu velicinu
font_size = st.number_input('Enter a font size: ', value = 11)


uploaded_file = st.file_uploader("Choose a file", type = ['ppt','pptx'])

if uploaded_file is not None:

    working_ppt = Presentation(uploaded_file)
    binary_output = BytesIO()

    nfc = RGBColor(0, 176, 80)

    slide_names = [f'Slide {i+1}' for i in range(len(working_ppt.slides))]

    select_all = st.checkbox('Apply to All', value = True)



    if select_all:
        opcija = True

    else:
        opcija = False
    
    # Create a list of checkboxes
    slide_checkboxes = []
    for i, slide in enumerate(working_ppt.slides):
        slide_name = f'slide {i+1}'
        slide_selected = st.checkbox(slide_name, value = opcija)
        if slide_selected:
            slide_checkboxes.append(True)
        else:
            slide_checkboxes.append(False)

 #--------------------umetnuti blok----------------------------------   
    #selected_slide_names = [name for name, selected in zip(slide_names, slide_selected) if selected]
    for slide, checkbox_value in zip(working_ppt.slides, slide_checkboxes):
        if checkbox_value:
            # Apply your formatting rules here
            shapes = slide.shapes
            for shape in shapes:
                if shape.shape_type == 19: #print(shape) for shape in shape.shape_type 19 je tabela:
                    table = shape.table
                    for cell in table.iter_cells():
                        paras = cell.text_frame.paragraphs
                        for para in paras:
                            if '%' not in cell.text:
                                for run in para.runs:
                                    nfn = selected_font       #input za font tip
                                    nfs = Pt(font_size)             #input za font size
                                    #nfc = RGBColor(0, 255, 0) input za font color 
                                    run.font.name = nfn
                                    run.font.size = nfs
                                    #run.font.color.rgb = nfc
                            elif '%' in cell.text:
                                res = "".join([ch for ch in cell.text if ch.isalpha()]) #(a,b,c) u abc
                                brojRegex = re.compile(r'\s(.*)') #brisanje svih posle razmaka
                                if brojRegex.search(cell.text): 
                                    try:
                                        mo = brojRegex.search(cell.text)
                                        paragraph = cell.text_frame.paragraphs[0]
                                        run = paragraph.add_run()
                                        run.text = brojRegex.sub(' ' + str(res),cell.text)
                                        cell.text_frame.clear()
                                        cell.text = run.text
                                        paragraph = cell.text_frame.paragraphs[0]
                                        paras = cell.text_frame.paragraphs
                                        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                                        for para in paras:
                                            for run in para.runs:
                                                nfn = selected_font        #input za font tip
                                                nfs = Pt(font_size)             #input za font size
                                                #nfc = RGBColor(0, 255, 0)#input za font color 
                                                run.font.name = nfn
                                                run.font.size = nfs
                                                run.font.color.rgb = nfc
                                                run.font.bold = True
                                            #print(cell.text)
                                    except AttributeError: #ako ne nadje
                                        continue
                                else:
                                    for run in para.runs:
                                        nfn = selected_font         #input za font tip
                                        nfs = Pt(font_size)             #input za font size
                                        #nfc = RGBColor(0, 255, 0)#input za font color 
                                        run.font.name = nfn
                                        run.font.size = nfs


            pass       



#-------------------KRAJ BLOKA-------------------------------
        

    # Display the selected slides
    #st.write(f'The selected slides are {selected_slides}')
    export_file_name = f"Formatted_ppt_{timestamp}.pptx"
    working_ppt.save(binary_output)
    export_ppt = working_ppt
    export_ppt.save(export_file_name)

    
    with open(export_file_name, 'rb') as f:
        file_bites = f.read()
        st.download_button(
            label = "Press to Download",
            data = binary_output.getvalue(),
            file_name = export_file_name,
            mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
else:
    st.info('☝️ You can just drag and drop it here')



#-------------TO DO----------------
#1. Dynamic chosing of style:
#    1.1. Choose font for all
#    1.2. Choose font size for all
#    1.3. Choose colors for affected cells
#
#2. Resolve text box formatting (headline and subheader issue)
#
#3. List all slides and ability to choose which slides they want
#formatting applied to
#----------------------------------

